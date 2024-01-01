import os

import pptx
import tkinter as tk
from pptx import Presentation
from pptx.dml.color import RGBColor
from datetime import datetime
from pptx.util import Inches, Pt
from tkinter import filedialog
pptx_path = pptx.__path__[0]


def list_files_and_folders(directory_path):
    result = {'files': [], 'folders': []}

    for root, dirs, files in os.walk(directory_path):
        current_folder = {'path': root, 'files': [], 'folders': []}

        for file in files:
            if file == "desktop.ini":
                continue
            current_folder['files'].append(file)

        for folder in dirs:
            current_folder['folders'].append(folder)

        result['files'].extend(current_folder['files'])
        result['folders'].append(current_folder)

    return result


def create_powerpoint_custom(directory_structure, output_pptx, max_items_per_slide=40):
    template_pptx_path = 'resources/Reporting-Bauakte-Template.pptx'
    prs = Presentation(template_pptx_path)

    # Use a default template, you can modify this to match your requirements
    title_slide_layout = prs.slide_layouts[0]
    folders_start_layout = prs.slide_layouts[1]
    default_folder_layout = prs.slide_layouts[2]
    docEstate_go_layout = prs.slide_layouts[3]

    # Add title slide
    title_slide = prs.slides.add_slide(title_slide_layout)

    # Add folders start slide with 2 small placeholders and 1 big placeholder
    folders_start_slide = prs.slides.add_slide(folders_start_layout)
    folders_start_shapes = folders_start_slide.shapes
    title_shape = folders_start_shapes.title
    right_big_placeholder = folders_start_shapes.placeholders[2]

    # Fill in the right big placeholder with the content of the first folder_info
    folder_info = directory_structure['folders'][0]
    title_shape.text = os.path.basename(folder_info['path'])

    for subfolder in folder_info['folders']:
        p = right_big_placeholder.text_frame.add_paragraph()
        p.text = f"Ordner: {subfolder}"
        p.font.size = Pt(11)  # Set font size for folders
        p.font.color.rgb = RGBColor(0, 142, 170)  # Set font color for folders

    for file in folder_info['files']:
        p = right_big_placeholder.text_frame.add_paragraph()
        p.text = f"Datei: {file}"
        p.font.size = Pt(9)  # Set font size for files
        p.font.color.rgb = RGBColor(66, 69, 72)  # Set font color for files

    # Loop through the rest of the folders using the default folder layout
    for folder_info in directory_structure['folders'][1:]:
        bullet_slide_layout = prs.slide_layouts[2]
        slide = prs.slides.add_slide(bullet_slide_layout)
        shapes = slide.shapes
        title_shape = shapes.title
        actual_right_placeholder = shapes.placeholders[2]
        actual_left_placeholder = shapes.placeholders[1]

        # Set the title of the slide to the current folder
        title_shape.text = f"Ordner: {os.path.basename(folder_info['path'])}"

        # List files in the folder
        items = [f"Ordner: {subfolder}" for subfolder in folder_info['folders']]

        # List subfolders in the folder
        items.extend([f"Datei: {file}" for file in folder_info['files']])

        # Add items to the slide with page breaks
        items_on_slide = 0
        for item in items:
            if items_on_slide % 2 == 0:
                # Use actual_left_placeholder for odd items (folders)
                p = actual_left_placeholder.text_frame.add_paragraph()
                p.text = item
                style_text(p, item)
            else:
                # Use actual_right_placeholder for even items (files)
                p = actual_right_placeholder.text_frame.add_paragraph()
                p.text = item
                style_text(p, item)

            items_on_slide += 1
            if items_on_slide >= max_items_per_slide:
                # Create a new slide if the maximum items per slide is reached
                slide = prs.slides.add_slide(default_folder_layout)
                shapes = slide.shapes
                title_shape = shapes.title
                actual_right_placeholder = shapes.placeholders[2]
                actual_left_placeholder = shapes.placeholders[1]
                items_on_slide = 0

    # Add the docEstate_go_layout slide
    prs.slides.add_slide(docEstate_go_layout)

    # Save the presentation
    prs.save(output_pptx)


def style_text(p, item):
    if item.startswith("Ordner"):
        p.font.size = Pt(11)  # Set font size for folders
        p.font.color.rgb = RGBColor(0, 142, 170)  # Set font color for folders
    else:
        p.font.size = Pt(9)  # Set font size for files
        p.font.color.rgb = RGBColor(66, 69, 72)  # Set font color for files


def get_directory(title):
    folder_selected = filedialog.askdirectory(title=title)
    return folder_selected


def main():
    # Create the main window
    root = tk.Tk()
    root.title("Generate power point for directory structure")

    # Get input and output directories
    input_folder = get_directory("Bitte Startverzeichnis auswählen")
    if len(input_folder) == 0:
        print(f"No folder selected -> abort")
        root.destroy()
        return
    output_folder = get_directory("Bitte Zielverzeichnis auswählen")

    # Get the directory structure
    directory_structure = list_files_and_folders(input_folder)

    # Create a PowerPoint presentation with a slide for each folder
    output_pptx = os.path.join(output_folder, f"Reporting-Bauakte {datetime.today().strftime('%d-%m-%Y')}.pptx")
    create_powerpoint_custom(directory_structure, output_pptx)

    print(f"PowerPoint file '{output_pptx}' created successfully.")

    tk.Label(root, text="Press any key to exit").pack()
    root.bind("<Key>", lambda event: root.destroy())
    root.mainloop()


if __name__ == "__main__":
    main()
